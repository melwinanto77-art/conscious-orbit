"""Document processor — extracts text from uploaded files for RAG indexing.

Falls back to empty string on any extraction failure so uploads are never
blocked by a corrupt or unsupported file.
"""
import csv
import logging

logger = logging.getLogger(__name__)

# MIME types we can extract text from. Images, zips and unknown types are
# skipped — uploads still succeed, they just won't be indexed for RAG.
PROCESSABLE_TYPES = {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "text/plain",
    "text/markdown",
    "text/csv",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/msword",
    "application/vnd.ms-excel",
    "application/vnd.ms-powerpoint",
}


def extract_text(file_path: str, content_type: str) -> str:
    """Extract plain text from a file based on its MIME type.

    Returns the extracted text, or empty string on failure.
    Never raises — extraction errors are logged and degraded gracefully.
    """
    try:
        if content_type == "application/pdf":
            return _extract_pdf(file_path)
        elif content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return _extract_docx(file_path)
        elif content_type in ("text/plain", "text/markdown"):
            return _extractPlainText(file_path)
        elif content_type == "text/csv":
            return _extract_csv(file_path)
        elif content_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            return _extract_xlsx(file_path)
        elif content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return _extract_pptx(file_path)
        elif content_type == "application/msword":
            return _extract_doc(file_path)
        elif content_type == "application/vnd.ms-excel":
            return _extract_xls(file_path)
        elif content_type == "application/vnd.ms-powerpoint":
            return _extract_ppt(file_path)
        else:
            return ""
    except Exception as e:
        logger.warning("Text extraction failed for %s (%s): %s", file_path, content_type, e)
        return ""


def _extract_pdf(path: str) -> str:
    from PyPDF2 import PdfReader
    reader = PdfReader(path)
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text)
    return "\n\n".join(pages)


def _extract_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extractPlainText(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def _extract_csv(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        rows = list(reader)
    if not rows:
        return ""
    lines = []
    for row in rows:
        lines.append(" | ".join(cell.strip() for cell in row))
    return "\n".join(lines)


def _extract_xlsx(path: str) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(sheets)


def _extract_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _extract_doc(path: str) -> str:
    """Legacy .doc — best-effort via python-docx (may fail on true .doc files)."""
    try:
        return _extract_docx(path)
    except Exception:
        return ""


def _extract_xls(path: str) -> str:
    """Legacy .xls — best-effort. Falls back to empty."""
    try:
        from openpyxl import load_workbook
        return _extract_xlsx(path)
    except Exception:
        return ""


def _extract_ppt(path: str) -> str:
    """Legacy .ppt — best-effort via python-pptx (may fail on true .ppt files)."""
    try:
        return _extract_pptx(path)
    except Exception:
        return ""


_text_splitter = None


def _get_splitter():
    global _text_splitter
    if _text_splitter is None:
        from langchain_text_splitters import RecursiveCharacterTextSplitter
        _text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            separators=["\n\n", "\n", ". ", " ", ""],
            length_function=len,
        )
    return _text_splitter


def chunk_text(text: str) -> list[str]:
    """Split text into chunks suitable for embedding.

    Returns a list of text chunks. Empty input returns an empty list.
    """
    if not text or not text.strip():
        return []
    splitter = _get_splitter()
    return [chunk for chunk in splitter.split_text(text) if chunk.strip()]
